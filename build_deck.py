"""
Build the Behavioral Baseline framework deck.
Matches the Cisco/Splunk O11y Days template:
  - 960×540 pt slides
  - Dark navy background: #0D274D
  - Cyan accent:          #00BCEB
  - White body text
  - Split-panel content slides: left navy panel, right white panel
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# ── Palette ─────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x27, 0x4D)
CYAN   = RGBColor(0x00, 0xBC, 0xEB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)   # light grey for right panel
MGRAY  = RGBColor(0xCC, 0xD6, 0xE0)   # medium grey for dividers
DKGRAY = RGBColor(0x33, 0x4E, 0x68)   # dark navy variant for boxes

# ── Slide dimensions ─────────────────────────────────────────────────────────
W = Inches(10)   # 960 pt ≈ 10 in
H = Inches(5.63) # 540 pt ≈ 5.63 in

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True,
                line_spacing=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as PPt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    run = p.add_run()
    run.text = text
    run.font.size     = Pt(font_size)
    run.font.bold     = bold
    run.font.color.rgb = color
    run.font.name     = "Calibri"
    return txb


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt as PPt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = PPt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return p


def add_bullet_box(slide, left, top, width, height,
                   bullets, font_size=13, color=WHITE,
                   title=None, title_color=CYAN):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(font_size + 1)
        run.font.bold  = True
        run.font.color.rgb = title_color
        run.font.name  = "Calibri"
    for b in bullets:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = b
        run.font.size  = Pt(font_size)
        run.font.bold  = False
        run.font.color.rgb = color
        run.font.name  = "Calibri"
    return txb


def footer(slide, text="© 2025 Cisco and/or its affiliates. All rights reserved.   Cisco Confidential"):
    add_textbox(slide,
                left=Inches(0.3), top=H - Inches(0.32),
                width=Inches(9.4), height=Inches(0.28),
                text=text, font_size=7, color=RGBColor(0xAA, 0xC4, 0xD8),
                align=PP_ALIGN.LEFT)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title / cover."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)

    # Cyan vertical accent bar on left
    add_rect(s, left=Inches(0), top=Inches(0),
             width=Inches(0.18), height=H, fill=CYAN)

    # Main title — centered
    add_textbox(s,
                left=Inches(0.45), top=Inches(1.4),
                width=Inches(9.2),  height=Inches(1.4),
                text="Behavioral Anomaly Detection Framework",
                font_size=44, bold=True, color=CYAN,
                align=PP_ALIGN.CENTER)

    add_textbox(s,
                left=Inches(0.45), top=Inches(2.75),
                width=Inches(9.2),  height=Inches(0.6),
                text="AI-Powered Behavioral Anomaly Detection for Splunk Observability APM",
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(s,
                left=Inches(0.45), top=Inches(3.5),
                width=Inches(9.2),  height=Inches(0.4),
                text="Marc Bui  |  Splunk",
                font_size=13, bold=False, color=RGBColor(0xAA, 0xC4, 0xD8),
                align=PP_ALIGN.CENTER)

    footer(s)
    return s


def slide_problem(prs):
    """Slide 2 — The problem standard alerting can't solve."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, left=Inches(0), top=Inches(0),
             width=Inches(0.18), height=H, fill=CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "The Gap in Standard Alerting",
                font_size=26, bold=True, color=CYAN)

    # Divider line
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.45),
                "APM AutoDetect covers error rate, latency, and request rate. But some failures leave no metric fingerprint.",
                font_size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

    # 4 problem boxes
    problems = [
        ("Service Disappears",
         "A service that always appeared in traces\nsuddenly vanishes — no error rate increase,\nno latency spike. Pure structural silence."),
        ("New Error Signature",
         "A brand new exception type fires for the\nfirst time. Threshold-based alerting requires\na baseline rate — there is none."),
        ("Call Path Changed",
         "A request now flows through a new service\nit never touched before. Could be a mis-deploy\nor a new dependency silently added."),
        ("DB Caller Goes Silent",
         "A service that always called the database\nstops doing so. No error, no latency —\njust behavioral silence."),
    ]

    box_w = Inches(2.15)
    box_h = Inches(2.8)
    gap   = Inches(0.12)
    start = Inches(0.45)

    for i, (title, body) in enumerate(problems):
        x = start + i * (box_w + gap)
        add_rect(s, x, Inches(1.55), box_w, box_h, DKGRAY)
        add_rect(s, x, Inches(1.55), box_w, Inches(0.06), CYAN)
        add_textbox(s, x + Inches(0.12), Inches(1.68), box_w - Inches(0.2), Inches(0.45),
                    title, font_size=13, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.12), Inches(2.18), box_w - Inches(0.2), Inches(2.0),
                    body, font_size=11, color=WHITE)

    footer(s)
    return s


def slide_solution_overview(prs):
    """Slide 3 — What the framework is."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "What the Behavioral Anomaly Detection Framework Does",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.42),
                "Augments Splunk APM AutoDetect with a structural/behavioral detection layer. "
                "No YAML. No alert rules. No thresholds to configure.",
                font_size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Left column — what it learns
    add_rect(s, Inches(0.45), Inches(1.55), Inches(4.35), Inches(3.55), DKGRAY)
    add_rect(s, Inches(0.45), Inches(1.55), Inches(4.35), Inches(0.06), CYAN)
    add_textbox(s, Inches(0.57), Inches(1.65), Inches(4.1), Inches(0.4),
                "What it learns (once, from live traffic)",
                font_size=13, bold=True, color=CYAN)
    learn_bullets = [
        "  \u2022  Every known service-to-service call path",
        "  \u2022  Every known error signature per service",
        "  \u2022  Normal span count ranges per path",
        "  \u2022  Which services always appear in which traces",
    ]
    add_bullet_box(s, Inches(0.57), Inches(2.12), Inches(4.1), Inches(2.8),
                   learn_bullets, font_size=12, color=WHITE)

    # Right column — what it detects
    add_rect(s, Inches(4.95), Inches(1.55), Inches(4.7), Inches(3.55), DKGRAY)
    add_rect(s, Inches(4.95), Inches(1.55), Inches(4.7), Inches(0.06), CYAN)
    add_textbox(s, Inches(5.07), Inches(1.65), Inches(4.45), Inches(0.4),
                "What it detects (every 5 minutes, autonomously)",
                font_size=13, bold=True, color=CYAN)
    detect_bullets = [
        "  \u2022  Service missing from traces (MISSING_SERVICE)",
        "  \u2022  Unknown call path appeared (NEW_FINGERPRINT)",
        "  \u2022  Brand new error type, first occurrence fires",
        "  \u2022  2+ tiers on same service \u2192 correlated alert (up to [Critical] MULTI_TIER)",
        "  \u2022  Anomaly correlated to a recent deploy",
        "  \u2022  Baseline self-heals after an incident resolves",
    ]
    add_bullet_box(s, Inches(5.07), Inches(2.12), Inches(4.45), Inches(2.8),
                   detect_bullets, font_size=12, color=WHITE)

    footer(s)
    return s


def slide_architecture(prs):
    """Slide 4 — How it works (detection pipeline)."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "How It Works",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    # Pipeline row: 4 boxes + arrows
    stages = [
        ("1. Learn",
         "Samples live traces\nBuilds structural\nfingerprints and\nerror signature\nbaselines from\nreal traffic"),
        ("2. Watch\n(every 5 min)",
         "Compares new traces\nagainst baseline\nEmits custom events\nto Splunk on any\ndeviation"),
        ("3. Correlate\n(every 5 min)",
         "Joins Tier 1 + 2 + 3\nevents by service\nMULTI_TIER = Critical\nAnnotates with\ndeployment context\nDowngrades severity\nif deploy-correlated"),
        ("4. Agent\n(on demand)",
         "Claude (AWS Bedrock)\nreads all signals\nReasons holistically\nOutputs: severity,\nroot cause,\nrecommended action"),
    ]

    bw = Inches(2.1)
    bh = Inches(3.5)
    gap = Inches(0.18)
    sx = Inches(0.45)
    sy = Inches(1.1)

    for i, (title, body) in enumerate(stages):
        x = sx + i * (bw + gap)
        add_rect(s, x, sy, bw, bh, DKGRAY)
        add_rect(s, x, sy, bw, Inches(0.06), CYAN)
        add_textbox(s, x + Inches(0.1), sy + Inches(0.08), bw - Inches(0.15), Inches(0.55),
                    title, font_size=13, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.1), sy + Inches(0.7), bw - Inches(0.15), Inches(2.7),
                    body, font_size=11, color=WHITE)
        # Arrow (except after last box)
        if i < len(stages) - 1:
            ax = x + bw + Inches(0.02)
            add_textbox(s, ax, sy + Inches(1.5), Inches(0.16), Inches(0.4),
                        "\u25b6", font_size=14, bold=True, color=CYAN)

    # Bottom note
    add_textbox(s, Inches(0.45), Inches(4.72), Inches(9.2), Inches(0.35),
                "Splunk APM AutoDetect covers Tier 1b/3/4 (error rate, latency, request rate) natively "
                "for all APM environments — no provisioning needed.",
                font_size=10, color=RGBColor(0x99, 0xBB, 0xCC))

    footer(s)
    return s


def slide_autodetect_relationship(prs):
    """Slide 5 — Relationship with APM AutoDetect."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Built on Top of APM AutoDetect",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    # Two columns
    # Left: AutoDetect (what it already covers)
    add_rect(s, Inches(0.45), Inches(1.05), Inches(4.3), Inches(3.85), DKGRAY)
    add_rect(s, Inches(0.45), Inches(1.05), Inches(4.3), Inches(0.06), MGRAY)
    add_textbox(s, Inches(0.57), Inches(1.13), Inches(4.05), Inches(0.45),
                "Splunk APM AutoDetect  (always on, built-in)",
                font_size=13, bold=True, color=RGBColor(0xCC, 0xDD, 0xEE))
    ad_items = [
        "  \u2022  Error rate spike per service",
        "  \u2022  p99 latency drift per service",
        "  \u2022  Request rate anomaly on ingress",
        "",
        "  Fires automatically for every APM environment.",
        "  No provisioning, no configuration required.",
    ]
    add_bullet_box(s, Inches(0.57), Inches(1.65), Inches(4.05), Inches(3.0),
                   ad_items, font_size=12, color=RGBColor(0xBB, 0xCC, 0xDD))

    # Right: Behavioral Baseline (what it adds)
    add_rect(s, Inches(4.9), Inches(1.05), Inches(4.75), Inches(3.85), DKGRAY)
    add_rect(s, Inches(4.9), Inches(1.05), Inches(4.75), Inches(0.06), CYAN)
    add_textbox(s, Inches(5.02), Inches(1.13), Inches(4.5), Inches(0.45),
                "Behavioral Anomaly Detection adds on top",
                font_size=13, bold=True, color=CYAN)
    bb_items = [
        "  \u2022  Structural trace path drift",
        "  \u2022  Service missing from traces",
        "  \u2022  New error signature, first occurrence",
        "  \u2022  Cross-tier correlation (Tier 1 + 2 + 3)",
        "  \u2022  Deployment-aware severity downgrade",
        "  \u2022  Auto-promotion of new patterns",
        "  \u2022  Self-healing baseline after incidents",
        "  \u2022  Claude-generated triage + runbook",
    ]
    add_bullet_box(s, Inches(5.02), Inches(1.65), Inches(4.5), Inches(3.0),
                   bb_items, font_size=12, color=WHITE)

    footer(s)
    return s


def slide_tiers(prs):
    """Slide — Detection tier model."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Detection Tier Model",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.38),
                "Every anomaly event is tagged with the tier that fired it. When two or more tiers fire on the same service, "
                "correlate.py emits a single high-confidence correlated event.",
                font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Tier boxes
    tiers = [
        ("Tier 1b / 3 / 4",
         "APM AutoDetect\n(built-in, always on)",
         "Error rate spike\np99 latency drift\nRequest rate anomaly\n\nFires on metrics.\nNo config required.",
         RGBColor(0x55, 0x77, 0x99),   # muted — not this framework
         RGBColor(0xAA, 0xBB, 0xCC)),
        ("Tier 2",
         "Trace Path Drift\n(trace_fingerprint.py)",
         "MISSING_SERVICE\n  \u2192 expected service absent from traces\n\nNEW_FINGERPRINT\n  \u2192 unknown call path appeared\n\nSpan count spike\n  \u2192 unexpected extra hops",
         CYAN, WHITE),
        ("Tier 3",
         "Error Signature Drift\n(error_fingerprint.py)",
         "NEW_ERROR_SIGNATURE\n  \u2192 first-ever occurrence of\n     a new exception type\n\nSIGNATURE_VANISHED\n  \u2192 known error pattern gone\n     (service may be silent)",
         CYAN, WHITE),
        ("Correlation",
         "Cross-Tier Join\n(correlate.py)",
         "TIER2_TIER3  \u2192 Major\n  trace + error on same service\n\nTIER1_TIER2 / TIER1_TIER3\n  \u2192 Major\n\nMULTI_TIER  \u2192 Critical\n  all 3 tiers on same service\n  = highest confidence",
         RGBColor(0xFF, 0xA5, 0x00), WHITE),  # amber — the output
    ]

    bw = Inches(2.1)
    bh = Inches(3.25)
    gap = Inches(0.18)
    sy  = Inches(1.5)

    for i, (tier_label, subtitle, body, bar_color, text_color) in enumerate(tiers):
        x = Inches(0.45) + i * (bw + gap)
        add_rect(s, x, sy, bw, bh, DKGRAY)
        add_rect(s, x, sy, bw, Inches(0.06), bar_color)
        add_textbox(s, x + Inches(0.12), sy + Inches(0.1), bw - Inches(0.2), Inches(0.38),
                    tier_label, font_size=13, bold=True, color=bar_color)
        add_textbox(s, x + Inches(0.12), sy + Inches(0.5), bw - Inches(0.2), Inches(0.38),
                    subtitle, font_size=10, bold=False, color=RGBColor(0xAA, 0xC4, 0xD8))
        add_textbox(s, x + Inches(0.12), sy + Inches(0.92), bw - Inches(0.2), Inches(2.2),
                    body, font_size=10, color=text_color)
        # Arrow between boxes (except after last)
        if i < len(tiers) - 1:
            ax = x + bw + Inches(0.02)
            add_textbox(s, ax, sy + Inches(1.4), Inches(0.16), Inches(0.4),
                        "\u25b6", font_size=14, bold=True, color=CYAN)

    # Bottom callout
    add_rect(s, Inches(0.45), Inches(4.88), Inches(9.2), Inches(0.38), RGBColor(0x1A, 0x3A, 0x5C))
    add_textbox(s, Inches(0.57), Inches(4.9), Inches(9.0), Inches(0.34),
                "Demo 4:  vets-service + DB down  \u2192  Tier 1 (AutoDetect error rate)  +  Tier 2 (MISSING_SERVICE)  +  Tier 3 (NEW_ERROR_SIGNATURE)  "
                "\u2192  correlate.py emits  [Critical] MULTI_TIER",
                font_size=10, bold=False, color=RGBColor(0xFF, 0xCC, 0x44))

    footer(s)
    return s


def slide_petclinic_topology(prs):
    """Slide — Demo environment: Spring PetClinic topology."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Demo Environment: Spring PetClinic on Kubernetes",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.38),
                "All demos run against a live Spring PetClinic deployment on k3d (k8s). "
                "Splunk OTel Java agent auto-instruments every service.",
                font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    BOX_W = Inches(2.0)
    BOX_H = Inches(0.56)
    INFRA  = RGBColor(0x33, 0x4E, 0x68)
    GW     = RGBColor(0x1A, 0x5C, 0x4A)
    SVC    = RGBColor(0x1A, 0x3A, 0x5C)
    DB     = RGBColor(0x5C, 0x2A, 0x1A)

    def svc_box(slide, x, y, label, sublabel, bar_color, fill_color):
        add_rect(slide, x, y, BOX_W, BOX_H, fill_color)
        add_rect(slide, x, y, BOX_W, Inches(0.05), bar_color)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.07),
                    BOX_W - Inches(0.15), Inches(0.28),
                    label, font_size=11, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.32),
                    BOX_W - Inches(0.15), Inches(0.22),
                    sublabel, font_size=9, bold=False, color=RGBColor(0xAA, 0xC4, 0xD8))

    # Row 1 — infra services
    row1_y = Inches(1.5)
    svc_box(s, Inches(1.5),  row1_y, "config-server",     "Spring Cloud Config", MGRAY, INFRA)
    svc_box(s, Inches(4.0),  row1_y, "discovery-server",  "Eureka Service Registry", MGRAY, INFRA)
    svc_box(s, Inches(6.5),  row1_y, "admin-server",      "Spring Boot Admin", MGRAY, INFRA)

    # Row 2 — gateway
    row2_y = Inches(2.25)
    svc_box(s, Inches(3.75), row2_y, "api-gateway",       "Public ingress / load balancer", CYAN, GW)

    # Row 3 — core services
    row3_y = Inches(3.0)
    svc_box(s, Inches(1.2),  row3_y, "customers-service", "Owner + pet profiles", CYAN, SVC)
    svc_box(s, Inches(3.75), row3_y, "vets-service",      "Veterinarian catalog", CYAN, SVC)
    svc_box(s, Inches(6.3),  row3_y, "visits-service",    "Appointment records", CYAN, SVC)

    # Row 4 — database
    row4_y = Inches(3.75)
    svc_box(s, Inches(3.75), row4_y, "petclinic-db",      "MySQL (shared data store)", RGBColor(0xFF, 0x66, 0x22), DB)

    # Arrows
    def arrow_down(slide, x, y):
        add_textbox(slide, x, y, Inches(0.2), Inches(0.25),
                    "\u25bc", font_size=9, bold=False, color=RGBColor(0x66, 0x99, 0xBB))

    arrow_down(s, Inches(4.75), row1_y + BOX_H + Inches(0.02))   # infra row gap
    arrow_down(s, Inches(4.75), row2_y + BOX_H + Inches(0.02))   # gw -> services
    arrow_down(s, Inches(2.15), row3_y + BOX_H + Inches(0.02))   # customers -> db
    arrow_down(s, Inches(4.75), row3_y + BOX_H + Inches(0.02))   # vets -> db
    arrow_down(s, Inches(7.25), row3_y + BOX_H + Inches(0.02))   # visits -> db

    # Legend bar
    add_rect(s, Inches(0.45), Inches(4.82), Inches(9.2), Inches(0.42), RGBColor(0x1A, 0x3A, 0x5C))
    add_textbox(s, Inches(0.57), Inches(4.84), Inches(9.0), Inches(0.38),
                "Loadgen hits api-gateway every ~5 s  \u00b7  OTel Java agent on all services  \u00b7  "
                "Traces + metrics flow to Splunk Observability (env: petclinicmbtest)",
                font_size=10, color=RGBColor(0xBB, 0xCC, 0xDD))

    footer(s)
    return s


def slide_demo_overview(prs):
    """Slide 6 — Demo agenda."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "What We'll Demo Today",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    demos = [
        ("Demo 0", "Steady State",
         "Framework in normal operation — baselines learned, cron running, zero anomalies"),
        ("Demo 1", "DB Outage \u2192 New Error Signatures",
         "Database goes down. New CannotCreateTransactionException fires on first occurrence. Claude: INCIDENT, PAGE_ONCALL"),
        ("Demo 2", "Bad Deploy \u2192 First-Occurrence Error",
         "visits-service crashes on startup. First request to the dead service triggers detection. No threshold exceeded"),
        ("Demo 3", "Missing Service \u2192 AI Triage",
         "vets-service killed. Framework detects structural absence from traces. Claude produces root cause + action in 3 min"),
        ("Demo 4", "All-Tier Correlation",
         "Both vets-service and DB down simultaneously. AutoDetect + trace drift + error signatures \u2192 [Critical] MULTI_TIER"),
        ("Demo 5", "Deploy-Correlated Severity Downgrade",
         "Bad deploy announced via CI/CD hook. correlate.py finds the deployment event and downgrades Major \u2192 Minor"),
        ("Demo 6", "Self-Healing",
         "New call path auto-promoted after 2 clean watch runs. Baseline healer scores windows and re-learns autonomously"),
        ("Demo 7", "Auto-Onboarding",
         "New environment discovered. Baselines built, dashboard created, cron scheduled, runbook generated \u2014 in 60 seconds"),
    ]

    col1_x = Inches(0.45)
    col2_x = Inches(1.6)
    col3_x = Inches(3.2)
    row_h  = Inches(0.535)
    start_y = Inches(1.06)

    for i, (num, title, desc) in enumerate(demos):
        y = start_y + i * row_h
        # Row alternating shade
        if i % 2 == 0:
            add_rect(s, Inches(0.45), y, Inches(9.2), row_h - Inches(0.03), DKGRAY)
        add_textbox(s, col1_x + Inches(0.08), y + Inches(0.05),
                    Inches(1.0), Inches(0.42),
                    num, font_size=11, bold=True, color=CYAN)
        add_textbox(s, col2_x, y + Inches(0.05),
                    Inches(1.5), Inches(0.42),
                    title, font_size=11, bold=True, color=WHITE)
        add_textbox(s, col3_x, y + Inches(0.05),
                    Inches(6.35), Inches(0.45),
                    desc, font_size=10, color=RGBColor(0xBB, 0xCC, 0xDD))

    # Demo note bar at bottom
    add_rect(s, Inches(0.45), H - Inches(0.55), Inches(9.2), Inches(0.38),
             RGBColor(0x1A, 0x3A, 0x5C))
    add_textbox(s, Inches(0.57), H - Inches(0.53), Inches(9.0), Inches(0.34),
                "\u26a0  Demo note: detections write to alerts.log + Splunk custom events.  "
                "Production target: full Splunk UI surfacing + Splunk On-Call / Slack notifications.",
                font_size=9, color=RGBColor(0xFF, 0xCC, 0x44))

    footer(s)
    return s


def slide_key_capabilities(prs):
    """Slide 7 — Key capabilities / value props for PM audience."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Key Capabilities",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    caps = [
        ("Zero Configuration",
         "No alert rules, no thresholds, no YAML.\nThe framework learns normal from live traffic\nand is ready in one command."),
        ("First-Occurrence Detection",
         "Fires the moment a new error signature or\nnew trace path appears — before any baseline\nrate exists to threshold against."),
        ("Claude-Powered Triage",
         "AWS Bedrock (Claude) synthesizes all signals\ninto a plain-English verdict: severity, root cause,\naffected services, recommended action."),
        ("Deployment-Aware",
         "Integrates with CI/CD via a one-line hook.\nAnomalies within 60 min of a deploy are\nauto-annotated and severity-downgraded."),
        ("Auto-Onboarding",
         "Runs every 30 min via cron. New environments\nare discovered, baselined, and fully operational\nwith no human intervention."),
        ("Self-Healing",
         "New patterns after deploys auto-promote after\n2 clean runs. Baselines re-learn autonomously\nafter an incident resolves."),
    ]

    bw = Inches(2.9)
    bh = Inches(2.0)
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)

    positions = [
        (Inches(0.45),          Inches(1.05)),
        (Inches(0.45) + bw + gap_x,   Inches(1.05)),
        (Inches(0.45) + 2*(bw + gap_x), Inches(1.05)),
        (Inches(0.45),          Inches(1.05) + bh + gap_y),
        (Inches(0.45) + bw + gap_x,   Inches(1.05) + bh + gap_y),
        (Inches(0.45) + 2*(bw + gap_x), Inches(1.05) + bh + gap_y),
    ]

    for (x, y), (title, body) in zip(positions, caps):
        add_rect(s, x, y, bw, bh, DKGRAY)
        add_rect(s, x, y, bw, Inches(0.06), CYAN)
        add_textbox(s, x + Inches(0.12), y + Inches(0.1), bw - Inches(0.2), Inches(0.4),
                    title, font_size=13, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.12), y + Inches(0.55), bw - Inches(0.2), Inches(1.35),
                    body, font_size=11, color=WHITE)

    footer(s)
    return s


def slide_section_break(prs, label="Live Demo"):
    """Full-bleed section break slide."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)

    # Bold left bar
    add_rect(s, Inches(0), Inches(0), Inches(0.35), H, CYAN)

    add_textbox(s, Inches(0.6), Inches(1.8), Inches(8.5), Inches(1.2),
                label, font_size=54, bold=True, color=CYAN,
                align=PP_ALIGN.LEFT)

    footer(s)
    return s


def slide_next_steps(prs):
    """Closing slide — product proposal for PM audience."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Product Proposal: Make This Native to Splunk Observability Cloud",
                font_size=22, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(0.98), Inches(9.2), Inches(0.38),
                "The proof of concept is working in production today. "
                "The ask: productize this as a native platform capability.",
                font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Three pillars
    steps = [
        ("Native Platform Integration",
         "Behavioral baseline learning built into\nAPM onboarding — a toggle, not a script.\n\n"
         "  \u2022  Detections surface alongside AutoDetect\n"
         "       alerts in the Splunk UI\n"
         "  \u2022  Same notification routing (Splunk On-Call,\n"
         "       Slack, webhook)\n"
         "  \u2022  Same muting, SLO wiring, and RBAC\n"
         "  \u2022  No external scripts or cron jobs"),
        ("AI Triage as a Product Feature",
         "Claude-generated verdict becomes the\n'Explain this alert' experience.\n\n"
         "  \u2022  Every INCIDENT gets a triage summary\n"
         "       attached automatically\n"
         "  \u2022  Root cause + affected services +\n"
         "       recommended action in plain English\n"
         "  \u2022  Generated runbook linked from the\n"
         "       alert detail view"),
        ("Platform Differentiation",
         "AutoDetect covers the metric layer.\n"
         "Every vendor covers the metric layer.\n\n"
         "  \u2022  Structural + behavioral detection catches\n"
         "       failures that leave no metric fingerprint\n"
         "  \u2022  First-occurrence detection — no threshold\n"
         "       to tune, fires on first event\n"
         "  \u2022  Self-healing baseline — zero ops overhead\n"
         "  \u2022  Defensible, differentiated capability"),
    ]

    bw = Inches(2.9)
    bh = Inches(3.45)
    gap = Inches(0.18)
    sy  = Inches(1.45)

    for i, (title, body) in enumerate(steps):
        x = Inches(0.45) + i * (bw + gap)
        add_rect(s, x, sy, bw, bh, DKGRAY)
        add_rect(s, x, sy, bw, Inches(0.06), CYAN)
        # Step number badge
        add_rect(s, x + Inches(0.12), sy + Inches(0.1),
                 Inches(0.32), Inches(0.32), CYAN)
        add_textbox(s, x + Inches(0.12), sy + Inches(0.08),
                    Inches(0.32), Inches(0.35),
                    str(i + 1), font_size=13, bold=True, color=NAVY,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.52), sy + Inches(0.08),
                    bw - Inches(0.65), Inches(0.42),
                    title, font_size=13, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.12), sy + Inches(0.58),
                    bw - Inches(0.22), Inches(2.75),
                    body, font_size=10.5, color=WHITE)

    # Bottom tagline
    add_textbox(s, Inches(0.45), Inches(4.98), Inches(9.2), Inches(0.35),
                "Questions?  \u2014  Marc Bui  |  Splunk Observability",
                font_size=12, bold=False,
                color=RGBColor(0xAA, 0xC4, 0xD8),
                align=PP_ALIGN.CENTER)

    footer(s)
    return s


# ── Build ────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_tiers(prs)
    slide_solution_overview(prs)
    slide_architecture(prs)
    slide_autodetect_relationship(prs)
    slide_key_capabilities(prs)
    slide_petclinic_topology(prs)
    slide_demo_overview(prs)
    slide_section_break(prs, "Live Demo")
    slide_next_steps(prs)

    out = "Behavioral_Baseline_Deck.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
