"""
Build the Behavioral Baseline framework deck.
Matches the Cisco/Splunk O11y Days template:
  - 960×540 pt slides
  - Dark navy background: #0D274D
  - Cyan accent:          #00BCEB
  - White body text
  - Split-panel content slides: left navy panel, right white panel
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# ── Palette ─────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x27, 0x4D)
CYAN   = RGBColor(0x00, 0xBC, 0xEB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)   # light grey for right panel
MGRAY  = RGBColor(0xCC, 0xD6, 0xE0)   # medium grey for dividers
DKGRAY = RGBColor(0x33, 0x4E, 0x68)   # dark navy variant for boxes

# ── Slide dimensions ─────────────────────────────────────────────────────────
W = Inches(10)   # 960 pt ≈ 10 in
H = Inches(5.63) # 540 pt ≈ 5.63 in

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True,
                line_spacing=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as PPt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    run = p.add_run()
    run.text = text
    run.font.size     = Pt(font_size)
    run.font.bold     = bold
    run.font.color.rgb = color
    run.font.name     = "Calibri"
    return txb


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt as PPt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = PPt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return p


def add_bullet_box(slide, left, top, width, height,
                   bullets, font_size=13, color=WHITE,
                   title=None, title_color=CYAN):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(font_size + 1)
        run.font.bold  = True
        run.font.color.rgb = title_color
        run.font.name  = "Calibri"
    for b in bullets:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = b
        run.font.size  = Pt(font_size)
        run.font.bold  = False
        run.font.color.rgb = color
        run.font.name  = "Calibri"
    return txb


def footer(slide, text="© 2025 Cisco and/or its affiliates. All rights reserved.   Cisco Confidential"):
    add_textbox(slide,
                left=Inches(0.3), top=H - Inches(0.32),
                width=Inches(9.4), height=Inches(0.28),
                text=text, font_size=7, color=RGBColor(0xAA, 0xC4, 0xD8),
                align=PP_ALIGN.LEFT)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title / cover."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)

    # Cyan vertical accent bar on left
    add_rect(s, left=Inches(0), top=Inches(0),
             width=Inches(0.18), height=H, fill=CYAN)

    # Main title — centered
    add_textbox(s,
                left=Inches(0.45), top=Inches(1.4),
                width=Inches(9.2),  height=Inches(1.4),
                text="Behavioral Anomaly Detection Framework",
                font_size=44, bold=True, color=CYAN,
                align=PP_ALIGN.CENTER)

    add_textbox(s,
                left=Inches(0.45), top=Inches(2.75),
                width=Inches(9.2),  height=Inches(0.6),
                text="AI-Powered Behavioral Anomaly Detection for Splunk Observability APM",
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(s,
                left=Inches(0.45), top=Inches(3.5),
                width=Inches(9.2),  height=Inches(0.4),
                text="Marc Bui  |  Splunk",
                font_size=13, bold=False, color=RGBColor(0xAA, 0xC4, 0xD8),
                align=PP_ALIGN.CENTER)

    footer(s)
    return s


def slide_problem(prs):
    """Slide 2 — The problem standard alerting can't solve."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, left=Inches(0), top=Inches(0),
             width=Inches(0.18), height=H, fill=CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "The Gap in Standard Alerting",
                font_size=26, bold=True, color=CYAN)

    # Divider line
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.45),
                "Standard metric alerting covers error rate, latency, and request rate. But some failures leave no metric fingerprint.",
                font_size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

    # 6 problem boxes (2 rows × 3)
    problems = [
        ("Service Disappears",
         "A service that always appeared in traces\nsuddenly vanishes — no error rate increase,\nno latency spike. Pure structural silence."),
        ("New Error Signature",
         "A brand new exception type fires for the\nfirst time. Threshold-based alerting requires\na baseline rate — there is none."),
        ("Call Path Changed",
         "A request now flows through a new service\nit never touched before. Could be a mis-deploy\nor a new dependency silently added."),
        ("DB Caller Goes Silent",
         "A service that always called the database\nstops doing so. No error, no latency —\njust behavioral silence."),
        ("Silent Span Collapse",
         "A service handles far fewer internal spans\nthan normal. No errors thrown, no latency\nchange — pipeline silently short-circuiting."),
        ("Span Count Explosion",
         "A service produces 10\u00d7 more spans than\nnormal — a retry storm or fan-out cascade\nthat errors won\u2019t surface on their own."),
    ]

    box_w = Inches(2.88)
    box_h = Inches(1.6)
    gap   = Inches(0.12)
    start = Inches(0.45)

    row_ys = [Inches(1.48), Inches(3.2)]
    for i, (title, body) in enumerate(problems):
        col = i % 3
        row = i // 3
        x = start + col * (box_w + gap)
        y = row_ys[row]
        add_rect(s, x, y, box_w, box_h, DKGRAY)
        add_rect(s, x, y, box_w, Inches(0.06), CYAN)
        add_textbox(s, x + Inches(0.12), y + Inches(0.1), box_w - Inches(0.2), Inches(0.42),
                    title, font_size=12, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.12), y + Inches(0.56), box_w - Inches(0.2), Inches(0.95),
                    body, font_size=10, color=WHITE)

    footer(s)
    return s


def slide_competitor_comparison(prs):
    """Slide — Competitive landscape: how others fill the behavioral detection gap."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9.2), Inches(0.55),
                "Competitors Are Already Filling This Gap",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.35),
                "Behavioral anomaly detection is a table-stakes capability in competitive APM platforms. "
                "Splunk O11y covers the metric layer — the structural layer is the gap.",
                font_size=11, color=RGBColor(0xCC, 0xDD, 0xEE))

    # ── Table layout ──────────────────────────────────────────────────────────
    # Columns: Capability | Dynatrace | Datadog | New Relic | Splunk (native) | This Framework
    col_labels  = ["", "Dynatrace\nSmartscape", "Datadog\nWatchdog", "New Relic\nApplied Intel.", "Splunk O11y\n(native)", "This\nFramework"]
    col_widths  = [Inches(2.52), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2)]
    col_x = [Inches(0.45)]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    # ✅ = native  ⚠ = partial  ❌ = not available
    CHECK  = "\u2705"   # ✅
    WARN   = "\u26a0\ufe0f"  # ⚠️  — use plain ⚠ for pptx compat
    CROSS  = "\u274c"   # ❌
    WARN   = "\u26a0"

    rows = [
        # (capability label, dynatrace, datadog, new relic, splunk native, this framework)
        ("Automatic topology\nmapping",              CHECK, CHECK, CHECK, CHECK, CHECK),
        ("Structural trace path\ndrift detection",   CHECK, WARN,  WARN,  CROSS, CHECK),
        ("Missing service\ndetection (~10s)",        CHECK, CROSS, CROSS, CROSS, CHECK),
        ("First-occurrence error\n(no threshold)",   WARN,  WARN,  WARN,  CROSS, CHECK),
        ("Cross-signal correlation\n(trace+err+metric)", CHECK, CHECK, WARN, CROSS, CHECK),
        ("Deploy-correlated\nseverity downgrade",    WARN,  WARN,  CROSS, CROSS, CHECK),
        ("AI root cause triage\n(LLM-generated)",    CHECK, CHECK, CHECK, CROSS, CHECK),
        ("Self-healing baseline\n(zero config)",     CHECK, WARN,  CROSS, CROSS, CHECK),
    ]

    ROW_H   = Inches(0.41)
    HDR_H   = Inches(0.38)
    tbl_top = Inches(1.42)

    # Header row
    for j, (label, w, x) in enumerate(zip(col_labels, col_widths, col_x)):
        bg = RGBColor(0x1A, 0x3A, 0x5C) if j == 0 else (DKGRAY if j < 5 else RGBColor(0x00, 0x6B, 0x85))
        add_rect(s, x, tbl_top, w, HDR_H, bg)
        txt_color = CYAN if j == len(col_labels) - 1 else RGBColor(0xCC, 0xDD, 0xEE)
        add_textbox(s, x + Inches(0.06), tbl_top + Inches(0.03),
                    w - Inches(0.1), HDR_H,
                    label, font_size=9, bold=True, color=txt_color,
                    align=PP_ALIGN.CENTER)

    # Data rows
    for i, row in enumerate(rows):
        y = tbl_top + HDR_H + i * ROW_H
        row_bg = RGBColor(0x12, 0x22, 0x3C) if i % 2 == 0 else RGBColor(0x0D, 0x1C, 0x32)
        for j, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
            if j == 0:
                # Capability label — left-aligned text
                add_rect(s, x, y, w, ROW_H, row_bg)
                add_textbox(s, x + Inches(0.1), y + Inches(0.05),
                            w - Inches(0.15), ROW_H,
                            cell, font_size=9.5, color=WHITE)
            else:
                # Vendor cell — icon centered, highlight last column (this framework)
                cell_bg = RGBColor(0x00, 0x4A, 0x5E) if j == len(row) - 1 else row_bg
                add_rect(s, x, y, w, ROW_H, cell_bg)
                icon_color = (CYAN if cell == CHECK else
                              RGBColor(0xFF, 0xCC, 0x44) if cell == WARN else
                              RGBColor(0xFF, 0x55, 0x55))
                add_textbox(s, x, y + Inches(0.05), w, ROW_H - Inches(0.05),
                            cell, font_size=13, bold=True, color=icon_color,
                            align=PP_ALIGN.CENTER)

    # Legend
    legend_y = tbl_top + HDR_H + len(rows) * ROW_H + Inches(0.06)
    add_textbox(s, Inches(0.45), legend_y, Inches(9.2), Inches(0.22),
                "\u2705 native capability    \u26a0 partial / requires config    \u274c not available",
                font_size=9, color=RGBColor(0x99, 0xBB, 0xCC),
                align=PP_ALIGN.LEFT)

    footer(s)
    return s


def slide_solution_overview(prs):
    """Slide 3 — What the framework is."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "What the Behavioral Anomaly Detection Framework Does",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.42),
                "An independent behavioral detection layer that runs inside the OTel Collector — "
                "no alert rules, no thresholds, no YAML.",
                font_size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Left column — what it learns
    add_rect(s, Inches(0.45), Inches(1.55), Inches(4.35), Inches(3.55), DKGRAY)
    add_rect(s, Inches(0.45), Inches(1.55), Inches(4.35), Inches(0.06), CYAN)
    add_textbox(s, Inches(0.57), Inches(1.65), Inches(4.1), Inches(0.4),
                "What it learns (once, from live traffic)",
                font_size=13, bold=True, color=CYAN)
    learn_bullets = [
        "  \u2022  Every known service-to-service call path",
        "  \u2022  Every known error signature per service",
        "  \u2022  Normal span count ranges per path",
        "  \u2022  Which services always appear in which traces",
    ]
    add_bullet_box(s, Inches(0.57), Inches(2.12), Inches(4.1), Inches(2.8),
                   learn_bullets, font_size=12, color=WHITE)

    # Right column — what it detects
    add_rect(s, Inches(4.95), Inches(1.55), Inches(4.7), Inches(3.55), DKGRAY)
    add_rect(s, Inches(4.95), Inches(1.55), Inches(4.7), Inches(0.06), CYAN)
    add_textbox(s, Inches(5.07), Inches(1.65), Inches(4.45), Inches(0.4),
                "What it detects (continuously, at the edge)",
                font_size=13, bold=True, color=CYAN)
    detect_bullets = [
        "  \u2022  Service missing from traces (MISSING_SERVICE)",
        "  \u2022  Unknown call path appeared (NEW_FINGERPRINT)",
        "  \u2022  Brand new error type, first occurrence fires",
        "  \u2022  New DB query template or slow query (z-score)",
        "  \u2022  Span count drop — silent failure, short-circuit",
        "  \u2022  Span count spike — retry storm, fan-out explosion",
        "  \u2022  Fires in ~10s — no poll interval, no Splunk wait",
        "  \u2022  2+ tiers on same service \u2192 correlated alert ([Critical] MULTI_TIER)",
        "  \u2022  Anomaly correlated to a recent deploy \u2192 severity downgrade",
        "  \u2022  Baseline self-heals after an incident resolves",
    ]
    add_bullet_box(s, Inches(5.07), Inches(2.12), Inches(4.45), Inches(2.8),
                   detect_bullets, font_size=12, color=WHITE)

    footer(s)
    return s


def slide_architecture(prs):
    """Slide 4 — How it works (detection pipeline)."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "How It Works",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    # Pipeline row: 4 boxes + arrows
    stages = [
        ("1. Learn",
         "Samples live traces\nBuilds structural\nfingerprints and\nerror signature\nbaselines from\nreal traffic"),
        ("2. Detect\n(~10s, at the edge)",
         "OTel Collector\nprocessor fingerprints\nevery trace inline\nas it flows through\n\nFires in ~10s —\nno poll interval,\nno Splunk wait"),
        ("3. Correlate\n(every 5 min)",
         "Joins Tier 1 + 2 + 3\nevents by service\nMULTI_TIER = Critical\nAnnotates with\ndeployment context\nDowngrades severity\nif deploy-correlated"),
        ("4. Agent\n(on demand)",
         "Claude (AWS Bedrock)\nreads all signals\nReasons holistically\nOutputs: severity,\nroot cause,\nrecommended action"),
    ]

    bw = Inches(2.1)
    bh = Inches(3.5)
    gap = Inches(0.18)
    sx = Inches(0.45)
    sy = Inches(1.1)

    for i, (title, body) in enumerate(stages):
        x = sx + i * (bw + gap)
        add_rect(s, x, sy, bw, bh, DKGRAY)
        add_rect(s, x, sy, bw, Inches(0.06), CYAN)
        add_textbox(s, x + Inches(0.1), sy + Inches(0.08), bw - Inches(0.15), Inches(0.55),
                    title, font_size=13, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.1), sy + Inches(0.7), bw - Inches(0.15), Inches(2.7),
                    body, font_size=11, color=WHITE)
        # Arrow (except after last box)
        if i < len(stages) - 1:
            ax = x + bw + Inches(0.02)
            add_textbox(s, ax, sy + Inches(1.5), Inches(0.16), Inches(0.4),
                        "\u25b6", font_size=14, bold=True, color=CYAN)

    # Bottom note
    add_textbox(s, Inches(0.45), Inches(4.72), Inches(9.2), Inches(0.35),
                "The OTel processor owns all behavioral signals — structural, metric, and error. "
                "Splunk APM receives the forwarded spans and provides the UI, dashboards, and storage.",
                font_size=10, color=RGBColor(0x99, 0xBB, 0xCC))

    footer(s)
    return s


def slide_detection_layers(prs):
    """Slide 5 — Two detection layers: OTel processor vs. Splunk APM."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Two Independent Detection Layers",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(0.97), Inches(9.2), Inches(0.38),
                "The OTel processor detects everything locally — structural, metric, and error signals — in ~10s. "
                "Splunk APM is the telemetry backend: it receives forwarded spans, stores traces, and provides the UI.",
                font_size=11.5, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Left: OTel Processor (this framework — owns all detection)
    add_rect(s, Inches(0.45), Inches(1.42), Inches(4.55), Inches(3.65), DKGRAY)
    add_rect(s, Inches(0.45), Inches(1.42), Inches(4.55), Inches(0.06), CYAN)
    add_textbox(s, Inches(0.57), Inches(1.5), Inches(4.3), Inches(0.45),
                "OTel Collector Processor  (this framework)",
                font_size=13, bold=True, color=CYAN)
    otel_items = [
        "  \u2022  MISSING_SERVICE — service absent from traces",
        "  \u2022  NEW_FINGERPRINT — unknown call path",
        "  \u2022  NEW_ERROR_SIGNATURE — first occurrence",
        "  \u2022  LATENCY_ANOMALY — z-score vs per-op baseline",
        "  \u2022  ERROR_RATE_ANOMALY — rate spike vs baseline",
        "  \u2022  SPAN_COUNT_DROP — silent failure / short-circuit",
        "  \u2022  SPAN_COUNT_SPIKE — retry storm / fan-out",
        "  \u2022  DB query: new template or slow query (z-score)",
        "",
        "  Fires in ~10s — no poll cycle, no indexing wait.",
        "  All baselines learned from live traffic.",
    ]
    add_bullet_box(s, Inches(0.57), Inches(2.02), Inches(4.3), Inches(2.9),
                   otel_items, font_size=11, color=WHITE)

    # Right: Splunk APM (telemetry backend + UI)
    add_rect(s, Inches(5.12), Inches(1.42), Inches(4.53), Inches(3.65), DKGRAY)
    add_rect(s, Inches(5.12), Inches(1.42), Inches(4.53), Inches(0.06), MGRAY)
    add_textbox(s, Inches(5.24), Inches(1.5), Inches(4.28), Inches(0.45),
                "Splunk Observability Cloud  (backend + UI)",
                font_size=13, bold=True, color=RGBColor(0xCC, 0xDD, 0xEE))
    splunk_items = [
        "  \u2022  Receives all forwarded spans (unchanged)",
        "  \u2022  Trace storage, search, and UI",
        "  \u2022  Service Map topology graph",
        "  \u2022  Dashboards and detector alerts",
        "  \u2022  Custom events from the processor",
        "       (trace.path.drift, error.signature.drift)",
        "",
        "  All detection is self-contained in the\n"
        "  OTel processor — no Splunk-side rules\n"
        "  or detectors required.",
    ]
    add_bullet_box(s, Inches(5.24), Inches(2.02), Inches(4.28), Inches(2.9),
                   splunk_items, font_size=11, color=RGBColor(0xBB, 0xCC, 0xDD))

    footer(s)
    return s


def slide_tiers(prs):
    """Slide — Detection tier model."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Detection Tier Model",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.38),
                "All detection runs inside the OTel Collector processor — structural, metric, and error signals in ~10s. "
                "correlate.py joins multi-tier signals into high-confidence correlated alerts.",
                font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Tier boxes
    tiers = [
        ("Tier 1",
         "Performance Anomalies\n(OTel processor — metric)",
         "LATENCY_ANOMALY\n  \u2192 z-score vs per-op baseline\n\nERROR_RATE_ANOMALY\n  \u2192 rate spike vs baseline\n\nSPAN_COUNT_DROP\n  \u2192 silent failure, short-circuit\n\nSPAN_COUNT_SPIKE\n  \u2192 retry storm, fan-out",
         RGBColor(0x00, 0x8C, 0xB0),   # teal — OTel metric signals
         WHITE),
        ("Tier 2",
         "Structural Drift\n(OTel processor — structural)",
         "MISSING_SERVICE\n  \u2192 expected service absent\n\nNEW_FINGERPRINT\n  \u2192 unknown call path appeared\n\nDB: new template or\n  slow query (z-score)",
         CYAN, WHITE),
        ("Tier 3",
         "Error Signature Drift\n(OTel processor — error)",
         "NEW_ERROR_SIGNATURE\n  \u2192 first-ever occurrence of\n     a new exception type\n\nSIGNATURE_VANISHED\n  \u2192 known error pattern gone\n     (service may be silent)",
         CYAN, WHITE),
        ("Correlation",
         "Cross-Tier Join\n(correlate.py)",
         "TIER1_TIER2  \u2192 Major\n  metric + structural\n\nTIER2_TIER3  \u2192 Major\n  structural + error\n\nMULTI_TIER  \u2192 Critical\n  all 3 tiers on same service\n  = highest confidence",
         RGBColor(0xFF, 0xA5, 0x00), WHITE),  # amber — the output
    ]

    bw = Inches(2.1)
    bh = Inches(3.25)
    gap = Inches(0.18)
    sy  = Inches(1.5)

    for i, (tier_label, subtitle, body, bar_color, text_color) in enumerate(tiers):
        x = Inches(0.45) + i * (bw + gap)
        add_rect(s, x, sy, bw, bh, DKGRAY)
        add_rect(s, x, sy, bw, Inches(0.06), bar_color)
        add_textbox(s, x + Inches(0.12), sy + Inches(0.1), bw - Inches(0.2), Inches(0.38),
                    tier_label, font_size=13, bold=True, color=bar_color)
        add_textbox(s, x + Inches(0.12), sy + Inches(0.5), bw - Inches(0.2), Inches(0.38),
                    subtitle, font_size=10, bold=False, color=RGBColor(0xAA, 0xC4, 0xD8))
        add_textbox(s, x + Inches(0.12), sy + Inches(0.92), bw - Inches(0.2), Inches(2.2),
                    body, font_size=10, color=text_color)
        # Arrow between boxes (except after last)
        if i < len(tiers) - 1:
            ax = x + bw + Inches(0.02)
            add_textbox(s, ax, sy + Inches(1.4), Inches(0.16), Inches(0.4),
                        "\u25b6", font_size=14, bold=True, color=CYAN)

    # Bottom callout
    add_rect(s, Inches(0.45), Inches(4.88), Inches(9.2), Inches(0.38), RGBColor(0x1A, 0x3A, 0x5C))
    add_textbox(s, Inches(0.57), Inches(4.9), Inches(9.0), Inches(0.34),
                "Demo 4:  vets-service + DB down  \u2192  Tier 3 fires in ~10s (error signatures)  +  Tier 2 fires at ~60s (MISSING_SERVICE)  "
                "\u2192  correlate.py emits  [Major] TIER2_TIER3",
                font_size=10, bold=False, color=RGBColor(0xFF, 0xCC, 0x44))

    footer(s)
    return s


def slide_petclinic_topology(prs):
    """Slide — Demo environment: Spring PetClinic topology."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Demo Environment: Spring PetClinic on Kubernetes",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(1.0), Inches(9.2), Inches(0.38),
                "All demos run against a live Spring PetClinic deployment on k3d (k8s). "
                "Splunk OTel Java agent auto-instruments every service.",
                font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    BOX_W = Inches(2.0)
    BOX_H = Inches(0.56)
    INFRA  = RGBColor(0x33, 0x4E, 0x68)
    GW     = RGBColor(0x1A, 0x5C, 0x4A)
    SVC    = RGBColor(0x1A, 0x3A, 0x5C)
    DB     = RGBColor(0x5C, 0x2A, 0x1A)

    def svc_box(slide, x, y, label, sublabel, bar_color, fill_color):
        add_rect(slide, x, y, BOX_W, BOX_H, fill_color)
        add_rect(slide, x, y, BOX_W, Inches(0.05), bar_color)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.07),
                    BOX_W - Inches(0.15), Inches(0.28),
                    label, font_size=11, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.32),
                    BOX_W - Inches(0.15), Inches(0.22),
                    sublabel, font_size=9, bold=False, color=RGBColor(0xAA, 0xC4, 0xD8))

    # Row 1 — infra services
    row1_y = Inches(1.5)
    svc_box(s, Inches(1.5),  row1_y, "config-server",     "Spring Cloud Config", MGRAY, INFRA)
    svc_box(s, Inches(4.0),  row1_y, "discovery-server",  "Eureka Service Registry", MGRAY, INFRA)
    svc_box(s, Inches(6.5),  row1_y, "admin-server",      "Spring Boot Admin", MGRAY, INFRA)

    # Row 2 — gateway
    row2_y = Inches(2.25)
    svc_box(s, Inches(3.75), row2_y, "api-gateway",       "Public ingress / load balancer", CYAN, GW)

    # Row 3 — core services
    row3_y = Inches(3.0)
    svc_box(s, Inches(1.2),  row3_y, "customers-service", "Owner + pet profiles", CYAN, SVC)
    svc_box(s, Inches(3.75), row3_y, "vets-service",      "Veterinarian catalog", CYAN, SVC)
    svc_box(s, Inches(6.3),  row3_y, "visits-service",    "Appointment records", CYAN, SVC)

    # Row 4 — database
    row4_y = Inches(3.75)
    svc_box(s, Inches(3.75), row4_y, "petclinic-db",      "MySQL (shared data store)", RGBColor(0xFF, 0x66, 0x22), DB)

    # Arrows
    def arrow_down(slide, x, y):
        add_textbox(slide, x, y, Inches(0.2), Inches(0.25),
                    "\u25bc", font_size=9, bold=False, color=RGBColor(0x66, 0x99, 0xBB))

    arrow_down(s, Inches(4.75), row1_y + BOX_H + Inches(0.02))   # infra row gap
    arrow_down(s, Inches(4.75), row2_y + BOX_H + Inches(0.02))   # gw -> services
    arrow_down(s, Inches(2.15), row3_y + BOX_H + Inches(0.02))   # customers -> db
    arrow_down(s, Inches(4.75), row3_y + BOX_H + Inches(0.02))   # vets -> db
    arrow_down(s, Inches(7.25), row3_y + BOX_H + Inches(0.02))   # visits -> db

    # Legend bar
    add_rect(s, Inches(0.45), Inches(4.82), Inches(9.2), Inches(0.42), RGBColor(0x1A, 0x3A, 0x5C))
    add_textbox(s, Inches(0.57), Inches(4.84), Inches(9.0), Inches(0.38),
                "Loadgen hits api-gateway every ~5 s  \u00b7  OTel Java agent on all services  \u00b7  "
                "Traces + metrics flow to Splunk Observability (env: test-7bb4-workshop)",
                font_size=10, color=RGBColor(0xBB, 0xCC, 0xDD))

    footer(s)
    return s


def slide_demo_overview(prs):
    """Slide 6 — Demo agenda."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "What We'll Demo Today",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    demos = [
        ("Demo 0", "Steady State",
         "Framework in normal operation — baselines learned, 0 drift events. 8-check pre-flight confirms every layer is ready."),
        ("Demo 1", "Kill Service \u2192 APM Still Green",
         "vets-service killed. APM Service Map shows green. OTel edge fires MISSING_SERVICE in ~10s. Claude detects before APM knows."),
        ("Demo 2", "New Call Path \u2192 Self-Healing",
         "New service added to call path. NEW_FINGERPRINT fires on runs 1+2, auto-promotes on run 2. Run 3: silence. Zero human intervention."),
        ("Demo 3", "New Error Signature \u2192 DB Kill",
         "DB goes down. CannotCreateTransactionException fires on first affected trace (~10s via OTel edge). Claude: INCIDENT."),
        ("Demo 4", "DB Gone Silent \u2192 Two Tiers",
         "vets-service + DB killed. Error tier fires in ~10s, MISSING_SERVICE at ~60s. Claude: mysql:petclinic root cause. TIER2_TIER3."),
        ("Demo 5", "Deploy-Correlated Severity Downgrade",
         "Bad deploy announced via notify_deployment.py. correlate.py finds the deployment event and downgrades Major \u2192 Minor [deployment-correlated]."),
        ("Demo 6", "Latency Spike",
         "visits-service latency spikes 250\u00d7 above baseline (z=8496). OTel edge detector fires LATENCY_ANOMALY in ~10s."),
        ("Demo 7", "Error Rate Anomaly",
         "customers-service error rate hits 100%. ERROR_RATE_ANOMALY + NEW_ERROR_SIGNATURE fire together. Claude: correlated incident."),
        ("Demo 8", "Combined Signal",
         "vets-service killed (structural) + customers-service error rate spikes (metric). Claude correlates all three tiers simultaneously."),
        ("Demo 9", "Slow DB",
         "DB overloaded — all callers get correlated LATENCY_ANOMALY spikes. Root cause: mysql:petclinic (shared dep, no structural drift)."),
        ("Demo 0b", "Auto-Onboarding",
         "New environment discovered. Baselines built, dashboard created, Claude runbook generated \u2014 in ~60 seconds."),
        ("Demo 11", "Span Count Drop (Topology UI)",
         "Connection pool exhaustion short-circuits DB calls on visits-service. No errors, no latency change. SPAN_COUNT_DROP fires in the UI (3 spans vs baseline 12\u201318). One-click Triage \u25b6 calls Claude."),
        ("Demo 12", "Span Count Spike (Topology UI)",
         "Retry storm on customers-service — 58 spans vs baseline max 12. SPAN_COUNT_SPIKE fires in the UI. Claude identifies fan-out cascade root cause."),
    ]

    col1_x = Inches(0.45)
    col2_x = Inches(1.6)
    col3_x = Inches(3.2)
    row_h  = Inches(0.415)
    start_y = Inches(1.06)

    for i, (num, title, desc) in enumerate(demos):
        y = start_y + i * row_h
        # Row alternating shade
        if i % 2 == 0:
            add_rect(s, Inches(0.45), y, Inches(9.2), row_h - Inches(0.03), DKGRAY)
        add_textbox(s, col1_x + Inches(0.08), y + Inches(0.04),
                    Inches(1.0), Inches(0.36),
                    num, font_size=10, bold=True, color=CYAN)
        add_textbox(s, col2_x, y + Inches(0.04),
                    Inches(1.5), Inches(0.36),
                    title, font_size=10, bold=True, color=WHITE)
        add_textbox(s, col3_x, y + Inches(0.04),
                    Inches(6.35), Inches(0.38),
                    desc, font_size=9, color=RGBColor(0xBB, 0xCC, 0xDD))

    # Demo note bar at bottom
    add_rect(s, Inches(0.45), H - Inches(0.55), Inches(9.2), Inches(0.38),
             RGBColor(0x1A, 0x3A, 0x5C))
    add_textbox(s, Inches(0.57), H - Inches(0.53), Inches(9.0), Inches(0.34),
                "\u26a0  Demo note: detections write to alerts.log + Splunk custom events.  "
                "Production target: full Splunk UI surfacing + Splunk On-Call / Slack notifications.",
                font_size=9, color=RGBColor(0xFF, 0xCC, 0x44))

    footer(s)
    return s


def slide_key_capabilities(prs):
    """Slide 7 — Key capabilities / value props for PM audience."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Key Capabilities",
                font_size=26, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    caps = [
        ("Zero Configuration",
         "No alert rules, no thresholds, no YAML.\nThe framework learns normal from live traffic\nand is ready in one command."),
        ("First-Occurrence Detection",
         "Fires the moment a new error signature or\nnew trace path appears — before any baseline\nrate exists to threshold against.\n\nAlso detects span count collapse and explosion\n— silent failures with no error or latency signal."),
        ("Claude-Powered Triage",
         "AWS Bedrock (Claude) synthesizes all signals\ninto a plain-English verdict: severity, root cause,\naffected services, recommended action."),
        ("Deployment-Aware",
         "Integrates with CI/CD via a one-line hook.\nAnomalies within 60 min of a deploy are\nauto-annotated and severity-downgraded."),
        ("Live Topology UI + RCA",
         "Real-time service graph with anomaly overlays\n(color-coded by signal type) and a one-click\nTriage \u25b6 button.\n\nClaude RCA panel synthesizes all active anomalies\ninto a structured root-cause assessment inline."),
        ("Self-Healing",
         "New patterns after deploys auto-promote after\n2 clean runs. Baselines re-learn autonomously\nafter an incident resolves."),
    ]

    bw = Inches(2.9)
    bh = Inches(2.0)
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)

    positions = [
        (Inches(0.45),          Inches(1.05)),
        (Inches(0.45) + bw + gap_x,   Inches(1.05)),
        (Inches(0.45) + 2*(bw + gap_x), Inches(1.05)),
        (Inches(0.45),          Inches(1.05) + bh + gap_y),
        (Inches(0.45) + bw + gap_x,   Inches(1.05) + bh + gap_y),
        (Inches(0.45) + 2*(bw + gap_x), Inches(1.05) + bh + gap_y),
    ]

    for (x, y), (title, body) in zip(positions, caps):
        add_rect(s, x, y, bw, bh, DKGRAY)
        add_rect(s, x, y, bw, Inches(0.06), CYAN)
        add_textbox(s, x + Inches(0.12), y + Inches(0.1), bw - Inches(0.2), Inches(0.4),
                    title, font_size=13, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.12), y + Inches(0.55), bw - Inches(0.2), Inches(1.35),
                    body, font_size=11, color=WHITE)

    footer(s)
    return s


def slide_section_break(prs, label="Live Demo"):
    """Full-bleed section break slide."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)

    # Bold left bar
    add_rect(s, Inches(0), Inches(0), Inches(0.35), H, CYAN)

    add_textbox(s, Inches(0.6), Inches(1.8), Inches(8.5), Inches(1.2),
                label, font_size=54, bold=True, color=CYAN,
                align=PP_ALIGN.LEFT)

    footer(s)
    return s


def slide_next_steps(prs):
    """Closing slide — product proposal for PM audience."""
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.45), Inches(0.25), Inches(9), Inches(0.55),
                "Product Proposal: Make This Native to Splunk Observability Cloud",
                font_size=22, bold=True, color=CYAN)
    add_rect(s, Inches(0.45), Inches(0.88), Inches(9.2), Inches(0.04), CYAN)

    add_textbox(s, Inches(0.45), Inches(0.98), Inches(9.2), Inches(0.38),
                "The proof of concept is working in production today. "
                "The ask: productize this as a native platform capability.",
                font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Three pillars
    steps = [
        ("Native Platform Integration",
         "Behavioral baseline learning built into\nAPM onboarding — a toggle, not a script.\n\n"
         "  \u2022  Detections surface natively in the\n"
         "       Splunk Observability UI\n"
         "  \u2022  Same notification routing (Splunk On-Call,\n"
         "       Slack, webhook)\n"
         "  \u2022  Same muting, SLO wiring, and RBAC\n"
         "  \u2022  No external scripts or cron jobs"),
        ("AI Triage as a Product Feature",
         "Claude-generated verdict becomes the\n'Explain this alert' experience.\n\n"
         "  \u2022  Every INCIDENT gets a triage summary\n"
         "       attached automatically\n"
         "  \u2022  Root cause + affected services +\n"
         "       recommended action in plain English\n"
         "  \u2022  Generated runbook linked from the\n"
         "       alert detail view"),
        ("Platform Differentiation",
         "Every APM vendor covers error rate, latency,\nand request rate. That\u2019s table stakes.\n\n"
         "  \u2022  Structural drift + silent failure detection\n"
         "       catches what metrics cannot surface\n"
         "  \u2022  First-occurrence detection — no threshold\n"
         "       to tune, fires on first event\n"
         "  \u2022  Self-healing baseline — zero ops overhead\n"
         "  \u2022  Defensible, differentiated capability"),
    ]

    bw = Inches(2.9)
    bh = Inches(3.45)
    gap = Inches(0.18)
    sy  = Inches(1.45)

    for i, (title, body) in enumerate(steps):
        x = Inches(0.45) + i * (bw + gap)
        add_rect(s, x, sy, bw, bh, DKGRAY)
        add_rect(s, x, sy, bw, Inches(0.06), CYAN)
        # Step number badge
        add_rect(s, x + Inches(0.12), sy + Inches(0.1),
                 Inches(0.32), Inches(0.32), CYAN)
        add_textbox(s, x + Inches(0.12), sy + Inches(0.08),
                    Inches(0.32), Inches(0.35),
                    str(i + 1), font_size=13, bold=True, color=NAVY,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.52), sy + Inches(0.08),
                    bw - Inches(0.65), Inches(0.42),
                    title, font_size=13, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.12), sy + Inches(0.58),
                    bw - Inches(0.22), Inches(2.75),
                    body, font_size=10.5, color=WHITE)

    # Bottom tagline
    add_textbox(s, Inches(0.45), Inches(4.98), Inches(9.2), Inches(0.35),
                "Questions?  \u2014  Marc Bui  |  Splunk Observability",
                font_size=12, bold=False,
                color=RGBColor(0xAA, 0xC4, 0xD8),
                align=PP_ALIGN.CENTER)

    footer(s)
    return s


# ── Build ────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_competitor_comparison(prs)
    slide_tiers(prs)
    slide_solution_overview(prs)
    slide_architecture(prs)
    slide_detection_layers(prs)
    slide_key_capabilities(prs)
    slide_petclinic_topology(prs)
    slide_demo_overview(prs)
    slide_section_break(prs, "Live Demo")
    slide_next_steps(prs)

    out = "Behavioral_Baseline_Deck.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
